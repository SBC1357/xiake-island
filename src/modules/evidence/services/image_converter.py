"""
多模态文件图像化转换服务

将 PDF / DOCX / DOC / PPTX / PPT / JPG / JPEG / PNG
统一转换为标准化的页面图像。

禁止 OCR 作为正式识别方案。
"""
import base64
import io
import hashlib
from pathlib import Path
from typing import List, Optional

from ..upload_models import (
    UploadedFile,
    PageImage,
    UploadFileType,
    generate_image_id,
)
from src.runtime_paths import get_upload_root


class ImageConversionError(Exception):
    pass


class ImageConverter:
    """
    文件到图像的统一转换器

    支持:
    - PDF -> 每页一张 PNG
    - DOCX/DOC -> 提取内容后生成文本快照图像
    - PPTX/PPT -> 提取幻灯片后生成快照图像
    - JPG/JPEG/PNG -> 直接作为标准化图像

    所有转换结果统一存储到运行态 uploads/{upload_id}/images/
    """

    def __init__(self, upload_root: Optional[Path] = None):
        self.upload_root = upload_root or get_upload_root()

    def convert(self, uploaded_file: UploadedFile) -> List[PageImage]:
        """
        将上传文件转换为页面图像列表

        Args:
            uploaded_file: 已上传的文件记录

        Returns:
            PageImage 列表
        """
        file_path = Path(uploaded_file.storage_path)
        if not file_path.exists():
            raise ImageConversionError(f"File not found: {file_path}")

        images_dir = self.upload_root / uploaded_file.upload_id / "images"
        images_dir.mkdir(parents=True, exist_ok=True)

        file_type = uploaded_file.file_type

        if file_type == UploadFileType.PDF:
            return self._convert_pdf(uploaded_file, file_path, images_dir)
        elif file_type in (UploadFileType.DOC, UploadFileType.DOCX):
            return self._convert_docx(uploaded_file, file_path, images_dir)
        elif file_type in (UploadFileType.PPT, UploadFileType.PPTX):
            return self._convert_pptx(uploaded_file, file_path, images_dir)
        elif file_type in (UploadFileType.JPG, UploadFileType.JPEG, UploadFileType.PNG):
            return self._convert_image(uploaded_file, file_path, images_dir)
        else:
            raise ImageConversionError(f"Unsupported file type: {file_type}")

    def _convert_pdf(
        self, uploaded_file: UploadedFile, file_path: Path, images_dir: Path
    ) -> List[PageImage]:
        """PDF -> 每页一张 PNG 图像"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImageConversionError(
                "PyMuPDF (fitz) is required for PDF conversion. "
                "Install: pip install PyMuPDF"
            )

        page_images = []
        doc = fitz.open(str(file_path))
        try:
            for page_num in range(len(doc)):
                page = doc[page_num]
                # 渲染为 2x 分辨率的 PNG
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)

                image_id = generate_image_id()
                image_path = images_dir / f"page_{page_num + 1}.png"
                pix.save(str(image_path))

                page_images.append(
                    PageImage(
                        image_id=image_id,
                        upload_id=uploaded_file.upload_id,
                        page_number=page_num + 1,
                        storage_path=str(image_path),
                        width=pix.width,
                        height=pix.height,
                        format="png",
                    )
                )
        finally:
            doc.close()

        return page_images

    def _convert_docx(
        self, uploaded_file: UploadedFile, file_path: Path, images_dir: Path
    ) -> List[PageImage]:
        """
        DOCX -> 文本快照图像

        当前策略：提取 DOCX 文本内容，使用 Pillow 生成文本快照图像。
        这保证了内容进入统一图像化主链。
        """
        try:
            from docx import Document
        except ImportError:
            raise ImageConversionError(
                "python-docx is required for DOCX conversion. "
                "Install: pip install python-docx"
            )

        doc = Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        if not paragraphs:
            paragraphs = ["[空文档]"]

        # 按每页约2000字符分页
        pages = self._split_text_to_pages(paragraphs, chars_per_page=2000)

        page_images = []
        for page_num, page_text in enumerate(pages, start=1):
            image_id = generate_image_id()
            image_path = images_dir / f"page_{page_num}.png"
            width, height = self._text_to_image(page_text, image_path)

            page_images.append(
                PageImage(
                    image_id=image_id,
                    upload_id=uploaded_file.upload_id,
                    page_number=page_num,
                    storage_path=str(image_path),
                    width=width,
                    height=height,
                    format="png",
                )
            )

        return page_images

    def _convert_pptx(
        self, uploaded_file: UploadedFile, file_path: Path, images_dir: Path
    ) -> List[PageImage]:
        """
        PPTX -> 每张幻灯片生成一张快照图像

        提取幻灯片文本和布局，使用 Pillow 渲染为图像。
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImageConversionError(
                "python-pptx is required for PPTX conversion. "
                "Install: pip install python-pptx"
            )

        prs = Presentation(str(file_path))
        page_images = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)

            slide_text = f"[幻灯片 {slide_num}]\n\n" + "\n".join(texts) if texts else f"[幻灯片 {slide_num} - 无文本内容]"

            image_id = generate_image_id()
            image_path = images_dir / f"slide_{slide_num}.png"
            width, height = self._text_to_image(slide_text, image_path)

            page_images.append(
                PageImage(
                    image_id=image_id,
                    upload_id=uploaded_file.upload_id,
                    page_number=slide_num,
                    storage_path=str(image_path),
                    width=width,
                    height=height,
                    format="png",
                )
            )

        return page_images

    def _convert_image(
        self, uploaded_file: UploadedFile, file_path: Path, images_dir: Path
    ) -> List[PageImage]:
        """JPG/JPEG/PNG -> 标准化为单页图像"""
        try:
            from PIL import Image
        except ImportError:
            raise ImageConversionError(
                "Pillow is required for image processing. "
                "Install: pip install Pillow"
            )

        img = Image.open(str(file_path))
        # 转换为 RGB（处理 RGBA）
        if img.mode == "RGBA":
            background = Image.new("RGB", img.size, (255, 255, 255))
            background.paste(img, mask=img.split()[3])
            img = background
        elif img.mode != "RGB":
            img = img.convert("RGB")

        image_id = generate_image_id()
        image_path = images_dir / f"page_1.png"
        img.save(str(image_path), "PNG")

        return [
            PageImage(
                image_id=image_id,
                upload_id=uploaded_file.upload_id,
                page_number=1,
                storage_path=str(image_path),
                width=img.width,
                height=img.height,
                format="png",
            )
        ]

    def _split_text_to_pages(
        self, paragraphs: List[str], chars_per_page: int = 2000
    ) -> List[str]:
        """将段落列表按字符数分页"""
        pages = []
        current_page = []
        current_chars = 0

        for para in paragraphs:
            if current_chars + len(para) > chars_per_page and current_page:
                pages.append("\n\n".join(current_page))
                current_page = []
                current_chars = 0
            current_page.append(para)
            current_chars += len(para)

        if current_page:
            pages.append("\n\n".join(current_page))

        return pages if pages else ["[空文档]"]

    def _text_to_image(self, text: str, output_path: Path) -> tuple[int, int]:
        """
        将文本渲染为 PNG 图像

        Returns:
            (width, height) 元组
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            raise ImageConversionError("Pillow is required. Install: pip install Pillow")

        # 设置画布参数
        width = 1200
        margin = 60
        line_height = 28
        font_size = 20

        # 尝试加载字体
        font = None
        try:
            font = ImageFont.truetype("msyh.ttc", font_size)  # 微软雅黑
        except (OSError, IOError):
            try:
                font = ImageFont.truetype("simsun.ttc", font_size)
            except (OSError, IOError):
                font = ImageFont.load_default()

        # 计算行数
        max_chars_per_line = (width - 2 * margin) // (font_size // 2 + 2)
        lines = []
        for paragraph in text.split("\n"):
            if not paragraph:
                lines.append("")
                continue
            while len(paragraph) > max_chars_per_line:
                lines.append(paragraph[:max_chars_per_line])
                paragraph = paragraph[max_chars_per_line:]
            lines.append(paragraph)

        height = max(400, 2 * margin + len(lines) * line_height + 40)

        # 创建图像
        img = Image.new("RGB", (width, height), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        y = margin
        for line in lines:
            draw.text((margin, y), line, fill=(0, 0, 0), font=font)
            y += line_height

        img.save(str(output_path), "PNG")
        return width, height


def image_to_base64(image_path: str) -> str:
    """将图像文件转换为 base64 字符串"""
    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")
